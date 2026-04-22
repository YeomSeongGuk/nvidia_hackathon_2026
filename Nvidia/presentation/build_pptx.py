"""Build hackathon pitch deck (5 min / 8 slides, English, NVIDIA theme).

Run:
    .venv/bin/python presentation/build_pptx.py

Output:
    presentation/hackathon_pitch.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- palette (NVIDIA-ish) ----
NVIDIA_GREEN = RGBColor(0x76, 0xB9, 0x00)
BG_DARK = RGBColor(0x0E, 0x1A, 0x2B)      # navy
BG_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x66, 0x70, 0x7A)
ACCENT_ORANGE = RGBColor(0xFF, 0x7A, 0x00)

# ---- layout ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]


def add_solid_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.insert(2, bg._element)  # send to back (pragma)


def add_green_bar(slide, y: Emu = Inches(0.5)) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(0.6)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NVIDIA_GREEN


def add_text(
    slide,
    text: str,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    font: str = "Helvetica",
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # margin inside
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    bullets: list[tuple[str, int]] | list[str],
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    *,
    size: int = 18,
    color: RGBColor = TEXT_DARK,
    muted: RGBColor = TEXT_MUTED,
    line_spacing: float = 1.2,
    font: str = "Helvetica",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, b in enumerate(bullets):
        if isinstance(b, tuple):
            text, level = b
        else:
            text, level = b, 0
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        bullet = "•" if level == 0 else "–"
        run.text = f"{bullet}  {text}"
        run.font.name = font
        run.font.size = Pt(size if level == 0 else size - 2)
        run.font.color.rgb = color if level == 0 else muted
        run.font.bold = level == 0
    return tb


def add_pill(slide, text: str, left: Emu, top: Emu, width: Emu, height: Emu,
             *, fill: RGBColor = NVIDIA_GREEN, text_color: RGBColor = BG_DARK, size: int = 14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Helvetica"
    return shape


def add_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    add_green_bar(slide, y=Inches(0.55))
    add_text(
        slide,
        title,
        Inches(0.8), Inches(0.45),
        Inches(12.0), Inches(0.75),
        size=32, bold=True, color=TEXT_DARK,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(0.8), Inches(1.15),
            Inches(12.0), Inches(0.45),
            size=16, color=TEXT_MUTED,
        )


# ============================================================
# Slide 1 — TITLE
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_DARK)

# NVIDIA green accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.2), Inches(2.2))
bar.fill.solid(); bar.fill.fore_color.rgb = NVIDIA_GREEN; bar.line.fill.background()

add_text(s, "Customer Driven Discovery",
         Inches(1.3), Inches(2.2), Inches(11), Inches(1.2),
         size=54, bold=True, color=TEXT_LIGHT)
add_text(s, "Building a Semantic Bridge from Reviews to Catalog",
         Inches(1.3), Inches(3.3), Inches(11), Inches(0.8),
         size=26, color=NVIDIA_GREEN)
add_text(s, "Korean fashion review curation powered by NeMo Curator + Nemotron-Super 120B",
         Inches(1.3), Inches(4.0), Inches(11), Inches(0.6),
         size=16, color=TEXT_MUTED)

# footer
add_text(s, "NVIDIA Hackathon 2026", Inches(0.8), Inches(6.8), Inches(6), Inches(0.4),
         size=12, color=TEXT_MUTED)
add_text(s, "Coupang Catalog Enrichment", Inches(7.0), Inches(6.8), Inches(5.5), Inches(0.4),
         size=12, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

# ============================================================
# Slide 2 — PROBLEM → SOLUTION
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_LIGHT)
add_slide_title(s, "The Gap: Intent vs Attributes",
                subtitle="Keyword search cannot reach the customer's real query")

# Problem box (left)
prob = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.8), Inches(1.9), Inches(5.6), Inches(2.3))
prob.fill.solid(); prob.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA); prob.line.fill.background()
add_text(s, "PROBLEM", Inches(1.0), Inches(2.0), Inches(5.0), Inches(0.4),
         size=13, bold=True, color=ACCENT_ORANGE)
add_text(
    s,
    '"내일 결혼식 갈 때 입을 옷" →\n'
    'Customer speaks in TPO / situation / emotion.\n\n'
    'Catalog stores only material / color / fit —\n'
    'a keyword engine finds nothing useful.',
    Inches(1.0), Inches(2.4), Inches(5.2), Inches(1.8),
    size=16, color=TEXT_DARK,
)

# Arrow
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           Inches(6.55), Inches(2.85), Inches(0.9), Inches(0.5))
arrow.fill.solid(); arrow.fill.fore_color.rgb = NVIDIA_GREEN; arrow.line.fill.background()

# Solution box (right)
sol = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(7.55), Inches(1.9), Inches(5.0), Inches(2.3))
sol.fill.solid(); sol.fill.fore_color.rgb = BG_DARK; sol.line.fill.background()
add_text(s, "SOLUTION", Inches(7.75), Inches(2.0), Inches(4.7), Inches(0.4),
         size=13, bold=True, color=NVIDIA_GREEN)
add_text(
    s,
    "Mine a Semantic Bridge from reviews.\n\n"
    "Reviews already carry TPO + attributes\n"
    "in natural language. Curate them with an\n"
    "LLM into canonical intents that map to\n"
    "catalog attributes.",
    Inches(7.75), Inches(2.4), Inches(4.7), Inches(1.8),
    size=16, color=TEXT_LIGHT,
)

# Bottom mapping strip
add_text(s, "reviews  →  canonical TPO intents  →  attribute bundles  →  natural-language queries",
         Inches(0.8), Inches(4.5), Inches(11.8), Inches(0.5),
         size=16, bold=True, color=NVIDIA_GREEN, align=PP_ALIGN.CENTER)

# Examples row
examples_y = Inches(5.2)
examples = [
    "일상용",
    "출근룩",
    "하객룩",
    "스크린골프",
    "필라테스 레이어",
]
ex_w = Inches(2.2); ex_gap = Inches(0.25)
total_w = ex_w * len(examples) + ex_gap * (len(examples) - 1)
start_x = (SLIDE_W - total_w) / 2
for i, ex in enumerate(examples):
    add_pill(s, ex, start_x + i * (ex_w + ex_gap), examples_y,
             ex_w, Inches(0.55), fill=NVIDIA_GREEN, size=15)

add_text(s, "260 canonical intents curated from 10 K Korean fashion reviews",
         Inches(0.8), Inches(6.2), Inches(11.8), Inches(0.5),
         size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 3 — NVIDIA + NeMo Stack  (3 × 2 grid)
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_LIGHT)
add_slide_title(s, "The NVIDIA + NeMo Stack",
                subtitle="Data pipeline · Korean-personalized synthesis · inference · independent eval")

# 3 × 2 grid
cards = [
    ("NeMo Curator",
     "Pipeline orchestration for Stage 1 — synth, filter, dedup, aggregate. Clean abstractions let every sub-stage iterate independently."),
    ("NeMo Data Designer",
     "Structured synthetic-data generation with seeded per-row LLM prompts. Stage 1.1 synthesizes reviews conditioned on seed + persona."),
    ("Korean Personalized Persona",
     "7-field persona space — age · sex · occupation · province · district · character · hobbies. Each synthetic review is anchored to a distinct Korean shopper."),
    ("Nemotron-Super 120B (FP8)",
     "Long-context Korean reasoning on Brev H100 NVL × 2, served via vLLM. enable_thinking=False for direct structured JSON."),
    ("Friendli Tri-Judge",
     "GLM-5.1 · DeepSeek-V3.2 · Qwen3-235B-A22B on serverless. Three independent foreign judges — Wisdom-of-the-Crowd eval."),
    ("BGE-M3 + scipy",
     "Korean semantic similarity drives Stage 2.2 agglomerative clustering and Stage 2.1.5 semantic dedup."),
]
card_w = Inches(3.95); card_h = Inches(2.25)
h_gap = Inches(0.2); v_gap = Inches(0.2)
total_row_w = card_w * 3 + h_gap * 2
start_x = (SLIDE_W - total_row_w) / 2
grid_top = Inches(1.9)

for i, (title, body) in enumerate(cards):
    row = i // 3
    col = i % 3
    x = start_x + col * (card_w + h_gap)
    y = grid_top + row * (card_h + v_gap)
    fill = BG_DARK if row == 0 else RGBColor(0x1A, 0x2B, 0x40)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid(); card.fill.fore_color.rgb = fill; card.line.fill.background()
    add_text(s, title,
             x + Inches(0.2), y + Inches(0.2), card_w - Inches(0.4), Inches(0.55),
             size=16, bold=True, color=NVIDIA_GREEN)
    add_text(s, body,
             x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), card_h - Inches(0.95),
             size=12, color=TEXT_LIGHT)

# footer row
add_text(s,
         "Deployed on Brev · H100 NVL × 2",
         Inches(0.8), Inches(6.75), Inches(11.8), Inches(0.4),
         size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 4 — WHY THIS STACK
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_LIGHT)
add_slide_title(s, "Why This Stack",
                subtitle="Domain fit, throughput, and evaluation independence")

rows = [
    ("Korean-first reasoning",
     "Nemotron-Super 120B handles Korean TPO nuance (존댓말/반말, 은어) without fine-tuning; few-shot + constrained JSON is enough."),
    ("Hackathon-scale throughput",
     "FP8 on H100 NVL × 2 keeps end-to-end latency low — 5 497 synthetic reviews curated to 260 canonicals in minutes, not hours."),
    ("NeMo Curator as pipeline spine",
     "Stage 1 synthetic + dedup + filter follow the Curator idioms. Clean abstraction boundaries let us iterate each stage independently."),
    ("Independent evaluation",
     "Friendli hosts 3 foreign judges — self-evaluation is a known failure mode; three diverse judges catch each other's blind spots."),
]
row_y = Inches(2.0)
for i, (k, v) in enumerate(rows):
    y = row_y + Inches(i * 1.1)
    add_text(s, "▸ " + k, Inches(0.9), y, Inches(4.2), Inches(0.6),
             size=18, bold=True, color=NVIDIA_GREEN)
    add_text(s, v, Inches(5.2), y, Inches(7.5), Inches(1.1),
             size=15, color=TEXT_DARK)

# ============================================================
# Slide 5 — APPROACH (Whole build + Eval + Iterate)
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_LIGHT)
add_slide_title(s, "Our Approach: Whole Build → Evaluate → Iterate",
                subtitle="Three-phase strategy — prove the end-to-end before tuning any single stage")

# 3 phase boxes
phase_w = Inches(4.0); phase_h = Inches(3.5); phase_y = Inches(2.0)
phase_gap = Inches(0.25)
phases_total_w = phase_w * 3 + phase_gap * 2
px = (SLIDE_W - phases_total_w) / 2

phases = [
    ("1 · WHOLE BUILD",
     "End-to-end pipeline first.",
     [
        "Stage 1: synth + dedup + filter",
        "Stage 2: extract → cluster → aggregate → expand",
        "10 K docs → 260 canonical TPOs → 1.3 K queries",
     ]),
    ("2 · EVAL INFRA",
     "Every stage gets its own judge.",
     [
        "5 LLM judge modules (stage_1_2, stage_2_1/2/3/4)",
        "3-way tri-judge ensemble via Friendli",
        "Deterministic probes for leak / format / diversity",
     ]),
    ("3 · ITERATE",
     "Hypothesis → patch → judge → promote.",
     [
        "One sealed subagent per iteration",
        "Eager-stacked promotion; rollback on regression",
        "Stage 1: 23 iters, best 8 / 10 gates",
     ]),
]
for i, (head, sub, bullets) in enumerate(phases):
    x = px + i * (phase_w + phase_gap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, phase_y, phase_w, phase_h)
    card.fill.solid(); card.fill.fore_color.rgb = BG_LIGHT; card.line.color.rgb = NVIDIA_GREEN
    card.line.width = Pt(2)
    add_text(s, head, x + Inches(0.25), phase_y + Inches(0.2), phase_w - Inches(0.5), Inches(0.5),
             size=16, bold=True, color=NVIDIA_GREEN)
    add_text(s, sub, x + Inches(0.25), phase_y + Inches(0.7), phase_w - Inches(0.5), Inches(0.45),
             size=13, color=TEXT_MUTED)
    add_bullets(s, bullets,
                x + Inches(0.25), phase_y + Inches(1.2),
                phase_w - Inches(0.5), phase_h - Inches(1.4),
                size=13, color=TEXT_DARK, line_spacing=1.3)

# Arrow connectors between cards
for i in range(2):
    x = px + (i + 1) * phase_w + i * phase_gap
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               x - Inches(0.05), phase_y + phase_h / 2 - Inches(0.15),
                               phase_gap + Inches(0.1), Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = NVIDIA_GREEN; arrow.line.fill.background()

# ============================================================
# Slide 6 — QUALITY ENGINE: Wisdom of the Crowd + LLM-as-Judge
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_LIGHT)
add_slide_title(s, "Quality Engine: Wisdom of the Crowd",
                subtitle="LLM-as-Judge × tri-judge consensus × deterministic probes")

# Left: 3 judge pills
jx = Inches(0.8); jy = Inches(2.0); jw = Inches(3.8); jh = Inches(3.8)
jcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, jx, jy, jw, jh)
jcard.fill.solid(); jcard.fill.fore_color.rgb = BG_DARK; jcard.line.fill.background()
add_text(s, "LLM-as-Judge × 3", jx + Inches(0.3), jy + Inches(0.2),
         jw - Inches(0.6), Inches(0.5),
         size=18, bold=True, color=NVIDIA_GREEN)
add_text(s, "Three independent foreign models\nvote on every stage output.",
         jx + Inches(0.3), jy + Inches(0.7), jw - Inches(0.6), Inches(0.9),
         size=13, color=TEXT_LIGHT)

judge_pills = [
    ("GLM-5.1", "Z.AI"),
    ("DeepSeek-V3.2", "DeepSeek"),
    ("Qwen3-235B-A22B", "Alibaba"),
]
for i, (name, vendor) in enumerate(judge_pills):
    top = jy + Inches(1.7) + Inches(i * 0.7)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              jx + Inches(0.3), top, jw - Inches(0.6), Inches(0.55))
    pill.fill.solid(); pill.fill.fore_color.rgb = NVIDIA_GREEN; pill.line.fill.background()
    add_text(s, f"{name}   ·   {vendor}",
             jx + Inches(0.35), top + Inches(0.05), jw - Inches(0.7), Inches(0.45),
             size=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

# Middle: formula
fx = Inches(4.9); fy = Inches(2.0); fw = Inches(3.5); fh = Inches(3.8)
fcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx, fy, fw, fh)
fcard.fill.solid(); fcard.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA); fcard.line.fill.background()
add_text(s, "Wisdom of the Crowd",
         fx + Inches(0.25), fy + Inches(0.2), fw - Inches(0.5), Inches(0.5),
         size=18, bold=True, color=ACCENT_ORANGE)
add_text(s,
         "• mean(judge₁, judge₂, judge₃)\n"
         "   → promote decision\n\n"
         "• range > 0.15\n"
         "   → HIGH_VARIANCE, block\n\n"
         "• agreement_rate on booleans\n"
         "   → confidence signal",
         fx + Inches(0.25), fy + Inches(0.9), fw - Inches(0.5), fh - Inches(1.0),
         size=14, color=TEXT_DARK)

# Right: deterministic probes
qx = Inches(8.7); qy = Inches(2.0); qw = Inches(3.9); qh = Inches(3.8)
qcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, qx, qy, qw, qh)
qcard.fill.solid(); qcard.fill.fore_color.rgb = BG_DARK; qcard.line.fill.background()
add_text(s, "Deterministic probes",
         qx + Inches(0.25), qy + Inches(0.2), qw - Inches(0.5), Inches(0.5),
         size=18, bold=True, color=NVIDIA_GREEN)
add_text(s,
         "Catch what judges miss:\n"
         "• title_reasoning_leak\n"
         "• dedup_miss_rate\n"
         "• canonical_non_fashion_rate\n"
         "• query_non_hangul_chars_rate\n"
         "• semdedup_retention\n\n"
         "Structured report per iter:\n"
         "judge_report · quant_report ·\n"
         "comparison · metrics.json",
         qx + Inches(0.25), qy + Inches(0.9), qw - Inches(0.5), qh - Inches(1.0),
         size=13, color=TEXT_LIGHT)

# Bottom line
add_text(s,
         "\"No single judge, no single metric — promote only when the crowd agrees and the probes confirm.\"",
         Inches(0.8), Inches(6.2), Inches(11.8), Inches(0.5),
         size=14, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 7 — RESULTS
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_LIGHT)
add_slide_title(s, "Results: 23 Iterations, Audited",
                subtitle="Stage 1 Δ on tri-judge mean — baseline vs iter_21_dedup_v2")

# Left: big stat strip
stats = [
    ("23", "iterations"),
    ("8 / 10", "promote gates (up from 1 / 10)"),
    ("260", "canonical TPOs from 5 497 synthetic reviews"),
    ("1.3 K", "natural-language queries generated"),
]
sx = Inches(0.8); sy = Inches(1.95); sw = Inches(5.0); sh = Inches(4.5)
sbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, sw, sh)
sbox.fill.solid(); sbox.fill.fore_color.rgb = BG_DARK; sbox.line.fill.background()
for i, (num, label) in enumerate(stats):
    y = sy + Inches(0.35) + Inches(i * 1.0)
    add_text(s, num, sx + Inches(0.3), y, Inches(2.0), Inches(0.9),
             size=44, bold=True, color=NVIDIA_GREEN)
    add_text(s, label, sx + Inches(2.4), y + Inches(0.15), sw - Inches(2.6), Inches(0.7),
             size=15, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)

# Right: delta table
tx = Inches(6.1); ty = Inches(1.95); tw = Inches(6.5); th = Inches(4.5)
tbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, ty, tw, th)
tbox.fill.solid(); tbox.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA); tbox.line.fill.background()
add_text(s, "Headline metric deltas", tx + Inches(0.25), ty + Inches(0.15),
         tw - Inches(0.5), Inches(0.5),
         size=17, bold=True, color=ACCENT_ORANGE)

deltas = [
    ("title_reasoning_leak",       "0.473",  "0.013"),
    ("fashion_rate",               "0.853",  "0.970"),
    ("persona_reflection (1-5)",   "3.62",   "4.66"),
    ("attr_grounded",              "0.467",  "0.990"),
    ("stage_1_2 retention",        "0.00",   "1.00"),
    ("rating_3_share",             "0.00",   "0.26"),
]
# header row
hy = ty + Inches(0.75)
add_text(s, "metric", tx + Inches(0.3), hy, Inches(3.3), Inches(0.4),
         size=13, bold=True, color=TEXT_MUTED)
add_text(s, "baseline", tx + Inches(3.6), hy, Inches(1.2), Inches(0.4),
         size=13, bold=True, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
add_text(s, "iter_21", tx + Inches(4.9), hy, Inches(1.4), Inches(0.4),
         size=13, bold=True, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
row_y0 = hy + Inches(0.45)
for i, (m, a, b) in enumerate(deltas):
    ry = row_y0 + Inches(i * 0.5)
    add_text(s, m, tx + Inches(0.3), ry, Inches(3.3), Inches(0.45), size=14, color=TEXT_DARK)
    add_text(s, a, tx + Inches(3.6), ry, Inches(1.2), Inches(0.45),
             size=14, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    add_text(s, "→  " + b, tx + Inches(4.9), ry, Inches(1.4), Inches(0.45),
             size=14, bold=True, color=NVIDIA_GREEN, align=PP_ALIGN.RIGHT)

# Footnote
add_text(s,
         "Remaining blockers: avg_text_quality (structural Naver floor) · dedup_miss_rate (semantic dedup queued)",
         Inches(0.8), Inches(6.6), Inches(11.8), Inches(0.4),
         size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 8 — HYPOTHESES: WINS & LOSSES
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_LIGHT)
add_slide_title(s, "Hypotheses: What Worked, What Didn't",
                subtitle="9 wins · 11 rejected · 2 deferred — each committed with a judge report")

# Column headers
wx = Inches(0.8); wy_hdr = Inches(1.95); wy = Inches(2.45)
ww = Inches(6.0)
lx = Inches(7.0); lw = Inches(5.5)

# WINS column
wbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, wx, wy_hdr - Inches(0.1),
                          ww, Inches(4.55))
wbox.fill.solid(); wbox.fill.fore_color.rgb = BG_DARK; wbox.line.fill.background()
add_pill(s, "WINS · kept", wx + Inches(0.25), wy_hdr, Inches(1.8), Inches(0.45),
         fill=NVIDIA_GREEN, size=13)
add_text(s, "9 hypotheses stacked into the best pipeline (iter_21)",
         wx + Inches(2.15), wy_hdr + Inches(0.05), ww - Inches(2.4), Inches(0.45),
         size=12, color=RGBColor(0xAA, 0xBB, 0xCC))

wins = [
    ("H11 · review_requires_attr_mention",
     "attr_grounded  0.467 → 0.99   (biggest single jump)"),
    ("H3 · title_postprocess",
     "title_leak  0.473 → 0.06   (best single-lever fix)"),
    ("H9 · korean_quality_filter",
     "stage_1_2 retention  0.00 → 1.00   (structural unblock)"),
    ("H12 · postgen_fashion_filter",
     "fashion_rate  0.867 → 0.953 → 0.977"),
    ("H8 · persona_binding_prompt",
     "persona_reflection  3.62 → 4.51;  rating-sentiment 0.78 → 1.0"),
    ("H4v2 · rating_category_sampler",
     "rating_3_share  0.00 → 0.30   (seed side couldn't — H4 falsified)"),
    ("H5v2 + H1 + H3v3  — stacked",
     "attr diversity · title prompt · leak regex   (small cumulative Δ)"),
]
row_y = wy_hdr + Inches(0.6)
for i, (head, body) in enumerate(wins):
    ry = row_y + Inches(i * 0.55)
    add_text(s, head, wx + Inches(0.3), ry, ww - Inches(0.6), Inches(0.3),
             size=13, bold=True, color=NVIDIA_GREEN)
    add_text(s, body, wx + Inches(0.3), ry + Inches(0.28), ww - Inches(0.6), Inches(0.3),
             size=11, color=TEXT_LIGHT)

# LOSSES column
lbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, wy_hdr - Inches(0.1),
                          lw, Inches(4.55))
lbox.fill.solid(); lbox.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA); lbox.line.fill.background()
add_pill(s, "ROLLED BACK", lx + Inches(0.25), wy_hdr, Inches(1.8), Inches(0.45),
         fill=ACCENT_ORANGE, size=13, text_color=BG_DARK)
add_text(s, "11 hypotheses falsified — each was a cheap negative signal",
         lx + Inches(2.15), wy_hdr + Inches(0.05), lw - Inches(2.4), Inches(0.45),
         size=12, color=TEXT_MUTED)

losses = [
    ("H4 · seed_rating_3_injection",
     "impossible — Naver corpus has 0 rating=3 rows"),
    ("H7 · non_fashion_seed_filter",
     "regression — exclude-list too aggressive, fashion_rate 0.85 → 0.79"),
    ("H3v2 · 45_char_cap",
     "regression — wider cap let reasoning fragments survive"),
    ("H3v4 · prose_reasoning_fallback",
     "regression — 4-char Korean fallback broke min-length format"),
    ("H14 / H14v2 · python_jaccard dedup",
     "no-op — pure-Python rule dedup fired 0 removals"),
    ("H2 · title_max_tokens_short",
     "partial — dominated by H3, not stacked"),
    ("H13 / H6 / H7v2",
     "mixed or trivially wrong; all committed as negative evidence"),
]
for i, (head, body) in enumerate(losses):
    ry = row_y + Inches(i * 0.55)
    add_text(s, head, lx + Inches(0.3), ry, lw - Inches(0.6), Inches(0.3),
             size=13, bold=True, color=ACCENT_ORANGE)
    add_text(s, body, lx + Inches(0.3), ry + Inches(0.28), lw - Inches(0.6), Inches(0.3),
             size=11, color=TEXT_DARK)

# Takeaway strip
add_text(s,
         "Falsifiability is cheap at 13 min per iter. Every \"loss\" was a ruled-out branch — "
         "not wasted work but a recorded negative signal in git.",
         Inches(0.8), Inches(6.7), Inches(11.8), Inches(0.5),
         size=13, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 9 — DATA + USE CASES + Q&A
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
add_solid_bg(s, BG_LIGHT)
add_slide_title(s, "Data + Use Cases",
                subtitle="What we curated and where it plugs in")

# Left: data acquired
lx = Inches(0.8); ly = Inches(1.95); lw = Inches(5.0); lh = Inches(4.5)
lbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ly, lw, lh)
lbox.fill.solid(); lbox.fill.fore_color.rgb = BG_DARK; lbox.line.fill.background()
add_text(s, "Data acquired", lx + Inches(0.3), ly + Inches(0.2),
         lw - Inches(0.6), Inches(0.5),
         size=18, bold=True, color=NVIDIA_GREEN)
add_bullets(
    s,
    [
        "10 K Korean fashion reviews (Musinsa, Naver Shopping, YouTube)",
        "Synthetic expansion → 5 497 curated reviews (Stage 1 pipeline)",
        "Curated to 260 canonical TPOs (Stage 2.2 clustering + LLM naming)",
        "Long-tail examples: 스크린골프, 낚시복, 필라테스 레이어, 분식룩",
        ("Head-tail coverage: 일상용 (1 145), 출근룩 (581), 데일리룩 (293)", 1),
    ],
    lx + Inches(0.3), ly + Inches(0.9), lw - Inches(0.6), lh - Inches(1.1),
    size=13, color=TEXT_LIGHT, muted=RGBColor(0xAA, 0xBB, 0xCC), line_spacing=1.35,
)

# Right: use cases
rx = Inches(6.1); ry = Inches(1.95); rw = Inches(6.5); rh = Inches(4.5)
rbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, ry, rw, rh)
rbox.fill.solid(); rbox.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA); rbox.line.fill.background()
add_text(s, "Use cases", rx + Inches(0.3), ry + Inches(0.2),
         rw - Inches(0.6), Inches(0.5),
         size=18, bold=True, color=ACCENT_ORANGE)
use_cases = [
    ("Semantic search",
     "Natural-language query → canonical TPO → attribute filter on catalog."),
    ("Recommendation",
     "Customer situation (upcoming event, weather, mood) → ranked intent bundle."),
    ("Ad / SEO copy",
     "1 canonical × 5 natural-language variants → ready-made keyword sets."),
    ("Catalog expansion",
     "Long-tail TPOs surface product gaps (e.g. 스크린골프 웨어 demand signal)."),
]
for i, (head, body) in enumerate(use_cases):
    uy = ry + Inches(1.0) + Inches(i * 0.88)
    add_text(s, "▸ " + head, rx + Inches(0.3), uy,
             Inches(2.5), Inches(0.45),
             size=14, bold=True, color=BG_DARK)
    add_text(s, body, rx + Inches(2.8), uy,
             rw - Inches(3.0), Inches(0.85),
             size=13, color=TEXT_DARK)

# Q&A strip
add_text(s, "Q & A", Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.5),
         size=20, bold=True, color=NVIDIA_GREEN, align=PP_ALIGN.CENTER)

# ============================================================
# write
# ============================================================
out_path = Path(__file__).resolve().parent / "hackathon_pitch.pptx"
out_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(out_path)
print(f"wrote {out_path}")
print(f"slides: {len(prs.slides)}")
